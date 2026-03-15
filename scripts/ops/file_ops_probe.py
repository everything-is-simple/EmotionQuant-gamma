from __future__ import annotations

import argparse
import importlib
import json
import os
import platform
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Probe local file-operation capabilities for Office, PDF, PS, PSD, and DuckDB."
    )
    parser.add_argument(
        "--temp-root",
        default=r"G:\EmotionQuant-temp\artifacts\file-ops-probe",
        help="Directory used for temporary probe artifacts.",
    )
    parser.add_argument(
        "--output",
        default=None,
        help="Optional JSON output path. Defaults to stdout only.",
    )
    return parser


def _module_status(name: str) -> dict[str, Any]:
    try:
        module = importlib.import_module(name)
    except Exception as exc:  # pragma: no cover - environment dependent
        return {
            "available": False,
            "error_type": type(exc).__name__,
            "error": str(exc),
        }
    return {
        "available": True,
        "version": getattr(module, "__version__", None),
    }


def _binary_status(name: str) -> dict[str, Any]:
    resolved = shutil.which(name)
    return {
        "available": resolved is not None,
        "path": resolved,
    }


def _probe_word(temp_root: Path) -> dict[str, Any]:
    from docx import Document

    path = temp_root / "probe.docx"
    document = Document()
    document.add_heading("EmotionQuant Word probe", level=1)
    document.add_paragraph("word-ready")
    document.save(path)

    reopened = Document(path)
    text = "\n".join(paragraph.text for paragraph in reopened.paragraphs)
    return {
        "available": "word-ready" in text,
        "path": str(path),
        "paragraph_count": len(reopened.paragraphs),
    }


def _probe_excel(temp_root: Path) -> dict[str, Any]:
    from openpyxl import Workbook, load_workbook

    path = temp_root / "probe.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Probe"
    worksheet["A1"] = "excel-ready"
    worksheet["B1"] = 42
    workbook.save(path)

    reopened = load_workbook(path, data_only=True)
    sheet = reopened["Probe"]
    return {
        "available": sheet["A1"].value == "excel-ready" and sheet["B1"].value == 42,
        "path": str(path),
        "sheet_names": reopened.sheetnames,
    }


def _probe_powerpoint(temp_root: Path) -> dict[str, Any]:
    from pptx import Presentation

    path = temp_root / "probe.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "EmotionQuant PPT probe"
    slide.placeholders[1].text = "ppt-ready"
    presentation.save(path)

    reopened = Presentation(path)
    title = reopened.slides[0].shapes.title.text
    return {
        "available": title == "EmotionQuant PPT probe",
        "path": str(path),
        "slide_count": len(reopened.slides),
    }


def _probe_pdf(temp_root: Path) -> dict[str, Any]:
    import fitz
    import pdfplumber
    from pypdf import PdfReader
    from reportlab.pdfgen import canvas

    path = temp_root / "probe.pdf"
    canvas_handle = canvas.Canvas(str(path))
    canvas_handle.drawString(72, 760, "pdf-ready")
    canvas_handle.save()

    reader_text = "".join(page.extract_text() or "" for page in PdfReader(str(path)).pages)
    with pdfplumber.open(str(path)) as plumber_doc:
        plumber_text = "".join(page.extract_text() or "" for page in plumber_doc.pages)
    with fitz.open(path) as pdf_doc:
        page_count = pdf_doc.page_count

    return {
        "available": "pdf-ready" in reader_text and "pdf-ready" in plumber_text,
        "path": str(path),
        "page_count": page_count,
    }


def _probe_duckdb(temp_root: Path) -> dict[str, Any]:
    import duckdb

    path = temp_root / "probe.duckdb"
    conn = duckdb.connect(str(path))
    try:
        conn.execute("create table if not exists probe(value varchar)")
        conn.execute("delete from probe")
        conn.execute("insert into probe values (?)", ["duckdb-ready"])
        value = conn.execute("select value from probe").fetchone()[0]
    finally:
        conn.close()

    return {
        "available": value == "duckdb-ready",
        "path": str(path),
    }


def _probe_postscript(temp_root: Path) -> dict[str, Any]:
    module_status = _module_status("ghostscript")
    binary_status = _binary_status("gswin64c")

    ps_path = temp_root / "probe.ps"
    pdf_path = temp_root / "probe-from-ps.pdf"
    ps_path.write_text(
        "%!PS-Adobe-3.0\n"
        "/Courier findfont 24 scalefont setfont\n"
        "72 720 moveto\n"
        "(EmotionQuant PS probe) show\n"
        "showpage\n",
        encoding="ascii",
    )

    if not binary_status["available"]:
        return {
            "available": False,
            "path": str(ps_path),
            "ghostscript_python": module_status,
            "ghostscript_binary": binary_status,
            "reason": "Ghostscript binary not found in PATH.",
        }

    command = [
        binary_status["path"],
        "-dNOPAUSE",
        "-dBATCH",
        "-sDEVICE=pdfwrite",
        f"-sOutputFile={pdf_path}",
        str(ps_path),
    ]
    completed = subprocess.run(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        check=False,
    )
    return {
        "available": completed.returncode == 0 and pdf_path.exists(),
        "path": str(ps_path),
        "rendered_pdf_path": str(pdf_path) if pdf_path.exists() else None,
        "ghostscript_python": module_status,
        "ghostscript_binary": binary_status,
        "returncode": completed.returncode,
        "stderr_tail": completed.stderr[-400:],
    }


def _probe_psd(temp_root: Path) -> dict[str, Any]:
    try:
        from psd_tools import PSDImage
    except Exception as exc:  # pragma: no cover - environment dependent
        return {
            "available": False,
            "error_type": type(exc).__name__,
            "error": str(exc),
        }

    path = temp_root / "probe.psd"
    psd = PSDImage.new(mode="RGB", size=(32, 32), color=(255, 255, 255))
    psd.save(path)
    reopened = PSDImage.open(path)
    return {
        "available": reopened.size == (32, 32),
        "path": str(path),
        "size": reopened.size,
    }


def _probe_com_app(progid: str, quit_method: str = "Quit") -> dict[str, Any]:
    win32_status = _module_status("win32com.client")
    if not win32_status["available"]:
        return {
            "available": False,
            "reason": "win32com.client is not available.",
        }

    try:
        code = (
            "import json, sys, win32com.client\n"
            "progid = sys.argv[1]\n"
            "quit_method = sys.argv[2]\n"
            "app = win32com.client.Dispatch(progid)\n"
            "try:\n"
            "    app.Visible = False\n"
            "except Exception:\n"
            "    pass\n"
            "try:\n"
            "    getattr(app, quit_method)()\n"
            "except Exception:\n"
            "    pass\n"
            "print(json.dumps({'available': True}))\n"
        )
        completed = subprocess.run(
            [sys.executable, "-c", code, progid, quit_method],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=15,
            check=False,
        )
        if completed.returncode != 0:
            return {
                "available": False,
                "error": completed.stderr.strip() or completed.stdout.strip(),
                "returncode": completed.returncode,
            }
        payload = json.loads(completed.stdout.strip() or "{}")
        return {"available": bool(payload.get("available"))}
    except subprocess.TimeoutExpired:
        return {
            "available": False,
            "error": "COM probe timed out after 15 seconds.",
        }
    except Exception as exc:  # pragma: no cover - environment dependent
        return {
            "available": False,
            "error_type": type(exc).__name__,
            "error": str(exc),
        }


def run_probe(temp_root: Path) -> dict[str, Any]:
    temp_root.mkdir(parents=True, exist_ok=True)

    payload: dict[str, Any] = {
        "python": {
            "version": sys.version,
            "executable": sys.executable,
            "platform": platform.platform(),
        },
        "modules": {
            name: _module_status(name)
            for name in [
                "docx",
                "pptx",
                "openpyxl",
                "pyxlsb",
                "msoffcrypto",
                "reportlab",
                "pdfplumber",
                "pypdf",
                "fitz",
                "duckdb",
                "win32com.client",
                "psd_tools",
                "ghostscript",
            ]
        },
        "binaries": {
            name: _binary_status(name)
            for name in ["WINWORD", "EXCEL", "POWERPNT", "Acrobat", "Photoshop", "gswin64c", "pdftoppm"]
        },
        "format_probes": {},
        "desktop_apps": {},
    }

    payload["format_probes"]["word_docx"] = _probe_word(temp_root)
    payload["format_probes"]["excel_xlsx"] = _probe_excel(temp_root)
    payload["format_probes"]["powerpoint_pptx"] = _probe_powerpoint(temp_root)
    payload["format_probes"]["pdf"] = _probe_pdf(temp_root)
    payload["format_probes"]["duckdb"] = _probe_duckdb(temp_root)
    payload["format_probes"]["postscript_ps"] = _probe_postscript(temp_root)
    payload["format_probes"]["photoshop_psd"] = _probe_psd(temp_root)

    if os.name == "nt":
        payload["desktop_apps"]["word_com"] = _probe_com_app("Word.Application")
        payload["desktop_apps"]["excel_com"] = _probe_com_app("Excel.Application")
        payload["desktop_apps"]["powerpoint_com"] = _probe_com_app("PowerPoint.Application")
        payload["desktop_apps"]["acrobat_com"] = _probe_com_app("AcroExch.App", quit_method="Exit")

    return payload


def main() -> int:
    args = build_parser().parse_args()
    temp_root = Path(args.temp_root).expanduser().resolve()
    payload = run_probe(temp_root)
    rendered = json.dumps(payload, ensure_ascii=False, indent=2)
    if args.output:
        output_path = Path(args.output).expanduser().resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(rendered, encoding="utf-8")
    print(rendered)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
